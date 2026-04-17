"""Génère presentation_lapage_v4.pptx — charte graphique personnalisée."""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import os

# ── Palette ───────────────────────────────────────────────────────────────────
GREEN       = RGBColor(0x1B, 0x43, 0x32)
GREEN_LITE  = RGBColor(0xA8, 0xD5, 0xBA)
KPI_GREEN   = RGBColor(0x1E, 0x8A, 0x44)
KPI_RED     = RGBColor(0xC0, 0x39, 0x2B)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x22, 0x22, 0x22)
MID         = RGBColor(0x77, 0x77, 0x77)
LIGHT_BG    = RGBColor(0xF2, 0xF2, 0xF2)
BORDER      = RGBColor(0xCC, 0xCC, 0xCC)

FONT        = "Aptos Display"
FIG         = "figures"

# ── Dimensions ────────────────────────────────────────────────────────────────
SLIDE_W     = Cm(33.87)
SLIDE_H     = Cm(19.05)
MARGIN      = Cm(2.0)           # marge gauche et droite
GAP         = Cm(0.4)           # gap graphique ↔ carrés
CONTENT_W   = SLIDE_W - 2 * MARGIN          # 29.87 cm
GRAPH_W     = (CONTENT_W - GAP) * 3 / 5    # ~17.68 cm
TEXT_W      = (CONTENT_W - GAP) * 2 / 5 - Cm(1.0)   # réduit de 1 cm
TEXT_X      = MARGIN + GRAPH_W + GAP
TITLE_Y          = Cm(0.45)
LISERET_GAP      = Cm(0.8)
CONTENT_GAP      = Cm(1.0)
BOTTOM_GRAPH     = Cm(1.0)
BOTTOM_INSIGHT   = Cm(1.0)
BOX_GAP          = Cm(0.2)
CARD_MARGIN      = Cm(2.0)   # marge gauche/droite des blocs colonnes (sommaire, synthèse, reco)
CARD_FONT        = 18        # taille police dans les cartes colonnes

BLANK = None

def init():
    global BLANK
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    BLANK = prs.slide_layouts[6]
    return prs


# ── Typographie ───────────────────────────────────────────────────────────────
def ft(run, size, bold=False, color=DARK, italic=False):
    run.font.name      = FONT
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.italic    = italic


# ── Fond blanc ────────────────────────────────────────────────────────────────
def white_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE


# ── Titre ─────────────────────────────────────────────────────────────────────
def add_title(slide, text):
    """Retourne le Y du bas du titre."""
    n_lines  = 1 if len(text) <= 55 else 2
    title_h  = Cm(2.8) if n_lines == 2 else Cm(1.6)
    tb = slide.shapes.add_textbox(MARGIN, TITLE_Y, CONTENT_W, title_h)
    tf = tb.text_frame
    tf.word_wrap = True
    r  = tf.paragraphs[0].add_run()
    r.text = text
    ft(r, 36, bold=True)
    return TITLE_Y + title_h


# ── Liseret ───────────────────────────────────────────────────────────────────
def add_liseret(slide, title_bottom_y):
    """Liseret pleine largeur du contenu. Retourne le Y du bas du liseret."""
    y = title_bottom_y + LISERET_GAP
    s = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        MARGIN, y, CONTENT_W, Cm(0.08)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = GREEN
    s.line.fill.background()
    return y + Cm(0.08)


# ── Numéro de slide ───────────────────────────────────────────────────────────
def add_num(slide, n):
    tb = slide.shapes.add_textbox(Cm(30.5), Cm(17.8), Cm(2.5), Cm(0.7))
    p  = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r  = p.add_run()
    r.text = str(n)
    ft(r, 10, color=MID)


# ── Header complet ─────────────────────────────────────────────────────────────
def header(slide, title, n):
    """Retourne content_y (Y de départ du contenu)."""
    white_bg(slide)
    title_bottom   = add_title(slide, title)
    liseret_bottom = add_liseret(slide, title_bottom)
    add_num(slide, n)
    return liseret_bottom + CONTENT_GAP


# ── Graphique ─────────────────────────────────────────────────────────────────
def add_figure(slide, filename, content_y):
    path = os.path.join(FIG, filename)
    if not os.path.exists(path):
        return
    h = SLIDE_H - content_y - BOTTOM_GRAPH
    slide.shapes.add_picture(path, MARGIN, content_y, GRAPH_W, h)


# ── Carrés d'insights ─────────────────────────────────────────────────────────
def add_insight_boxes(slide, groups, content_y):
    """
    groups : liste de dicts avec les clés :
      kpi        (opt) : chiffre clé affiché en grand
      kpi_color  (opt) : "green" | "red"  (défaut neutral)
      sub        (opt) : sous-label du kpi
      points     (opt) : liste de bullets
    """
    n          = len(groups)
    total_h    = SLIDE_H - content_y - BOTTOM_INSIGHT
    box_h      = (total_h - BOX_GAP * (n - 1)) / n

    for i, grp in enumerate(groups):
        y = content_y + i * (box_h + BOX_GAP)

        tb = slide.shapes.add_textbox(TEXT_X, y, TEXT_W, box_h)
        tb.fill.solid()
        tb.fill.fore_color.rgb = LIGHT_BG
        tb.line.color.rgb      = BORDER
        tb.line.width          = Pt(0.5)

        tf = tb.text_frame
        tf.word_wrap    = True
        tf.margin_left  = Cm(0.5)
        tf.margin_right = Cm(0.5)
        tf.margin_top   = Cm(0.4)
        tf.margin_bottom = Cm(0.3)

        first_p = True

        def next_para():
            nonlocal first_p
            if first_p:
                first_p = False
                return tf.paragraphs[0]
            p = tf.add_paragraph()
            p.space_before = Pt(4)
            return p

        # KPI
        kpi = grp.get("kpi")
        if kpi:
            color_name = grp.get("kpi_color", "neutral")
            kpi_color  = KPI_GREEN if color_name == "green" else \
                         KPI_RED   if color_name == "red"   else DARK
            p = next_para()
            r = p.add_run()
            r.text = kpi
            ft(r, 40, bold=True, color=kpi_color)

        # Sous-label du KPI
        sub = grp.get("sub")
        if sub:
            p = next_para()
            p.space_before = Pt(0)
            r = p.add_run()
            r.text = sub
            ft(r, 12, color=MID)

        # Bullets
        for bullet in grp.get("points", []):
            p = next_para()
            r = p.add_run()
            r.text = bullet
            ft(r, 14)


# ── Slide contenu standard ────────────────────────────────────────────────────
def content_slide(prs, title, fig_file, groups, n):
    sl = prs.slides.add_slide(BLANK)
    cy = header(sl, title, n)
    add_figure(sl, fig_file, cy)
    add_insight_boxes(sl, groups, cy)
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════
prs = init()


# ── Slide 1 — Couverture ─────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg = s1.background.fill
bg.solid()
bg.fore_color.rgb = GREEN

tb = s1.shapes.add_textbox(Cm(4), Cm(4.5), Cm(26), Cm(4.5))
tf = tb.text_frame
tf.word_wrap = True
p  = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r  = p.add_run()
r.text = "Lapage — Analyse e-commerce"
ft(r, 50, bold=True, color=WHITE)

deco = s1.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(8), Cm(9.6), Cm(18), Cm(0.07)
)
deco.fill.solid()
deco.fill.fore_color.rgb = GREEN_LITE
deco.line.fill.background()

tb2 = s1.shapes.add_textbox(Cm(4), Cm(10.2), Cm(26), Cm(2.5))
tf2 = tb2.text_frame
p2  = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
r2  = p2.add_run()
r2.text = "Mars 2021 – Février 2023  ·  679 111 transactions  ·  8 596 clients"
ft(r2, 36, color=GREEN_LITE)

tb3 = s1.shapes.add_textbox(Cm(4), Cm(15.5), Cm(26), Cm(1.5))
tf3 = tb3.text_frame
p3  = tf3.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER
r3  = p3.add_run()
r3.text = "Analyse Data — Avril 2026"
ft(r3, 16, color=RGBColor(0x7F, 0xBF, 0x9E))


# ── Slide 2 — Sommaire ───────────────────────────────────────────────────────
s2   = prs.slides.add_slide(BLANK)
cy2  = header(s2, "Sommaire", 2)

sections = [
    ("01", "Activité",     "CA, transactions\ncatalogue, produits"),
    ("02", "Clients",      "BtoB, concentration\nnouveaux clients"),
    ("03", "Corrélations", "Âge, genre\ncatégorie d'achat"),
    ("04", "Synthèse",     "Enseignements clés\nrecommandations"),
]
n_cols    = len(sections)
cards_w   = SLIDE_W - 2 * CARD_MARGIN
col_w     = (cards_w - Cm(0.3) * (n_cols - 1)) / n_cols
cards_h   = SLIDE_H - cy2 - Cm(2.0)
for i, (num, titre, desc) in enumerate(sections):
    x = CARD_MARGIN + i * (col_w + Cm(0.3))
    box = s2.shapes.add_textbox(x, cy2, col_w, cards_h)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = BORDER
    box.line.width = Pt(0.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Cm(0.45)

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = num
    ft(r1, 30, bold=True, color=GREEN)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(12)
    r2 = p2.add_run()
    r2.text = titre
    ft(r2, CARD_FONT, bold=True)

    p3 = tf.add_paragraph()
    p3.space_before = Pt(8)
    r3 = p3.add_run()
    r3.text = desc
    ft(r3, CARD_FONT, color=MID)


# ── Slides 3-11 — Activité ───────────────────────────────────────────────────
content_slide(prs,
    "CA stable sur 24 mois — pas de décrochage structurel, une base commerciale solide",
    "01_ca_mensuel.png",
    [
        {"kpi": "11,85 M€", "kpi_color": "green", "sub": "CA total · Mar 2021 – Fév 2023",
         "points": ["Mois record : Fév. 2022", "Mois creux : Oct. 2021"]},
        {"kpi": "493 k€", "kpi_color": "neutral", "sub": "CA moyen / mois",
         "points": ["Pas de tendance baissière sur 2 ans",
                    "→ L'activité est structurellement saine"]},
    ], 3)

content_slide(prs,
    "Cat. 1 et Cat. 0 dominent à 76% — Cat. 2 reste minoritaire malgré ses prix plus élevés",
    "07_repartition_categorie_pie.png",
    [
        {"kpi": "39 %", "kpi_color": "green", "sub": "Catégorie 1 — leader",
         "points": ["Catégorie 0 : 37 %", "Les deux ensemble : 76 % du CA"]},
        {"kpi": "23 %", "kpi_color": "red", "sub": "Catégorie 2 — minoritaire",
         "points": ["→ Offre, prix ou visibilité à investiguer",
                    "→ Potentiel de croissance non exploité"]},
    ], 4)

content_slide(prs,
    "Les 3 catégories évoluent en parallèle — pas de substitution, un marché cohérent",
    "02_ca_par_categorie.png",
    [
        {"points": ["Cat. 1 et Cat. 0 suivent la même saisonnalité",
                    "Aucune catégorie ne cannibalise les autres"]},
        {"points": ["→ Les promotions sur une catégorie n'affectent pas les autres",
                    "→ La croissance est globale, pas de rééquilibrage interne"]},
    ], 5)

content_slide(prs,
    "Cat. 2 a les prix unitaires les plus élevés — un positionnement premium sous-exploité",
    "07b_prix_unitaire_categ.png",
    [
        {"points": ["Cat. 2 : prix unitaire moyen nettement supérieur aux Cat. 0 et 1",
                    "→ Explique pourquoi Cat. 2 fait 23 % du CA avec peu de volume"]},
        {"points": ["→ Levier : cibler les tranches d'âge qui achètent Cat. 2 (< 25 ans)",
                    "→ Une réduction sur Cat. 2 a plus d'impact en € qu'une réduction sur Cat. 0"]},
    ], 6)

content_slide(prs,
    "Base clients stable sur 24 mois — l'acquisition compense le churn naturel",
    "03_clients_par_mois.png",
    [
        {"kpi": "5 714", "kpi_color": "neutral", "sub": "clients actifs en moyenne / mois",
         "points": ["Pas d'érosion visible sur la période", "La base est saine et constante"]},
        {"points": ["→ L'enjeu n'est pas l'acquisition brute mais la fidélisation",
                    "→ Travailler la valeur par client plutôt que le volume"]},
    ], 7)

content_slide(prs,
    "Cat. 0 et Cat. 1 mobilisent le plus de clients — Cat. 2 reste confidentielle mais premium",
    "07c_clients_categ_mois.png",
    [
        {"points": ["Répartition clients par catégorie stable sur toute la période",
                    "Cat. 2 : peu de clients mais panier unitaire élevé"]},
        {"points": ["→ Cat. 2 est un produit de niche — cibler les profils à fort potentiel",
                    "→ Ne pas chercher à massifier Cat. 2 : son modèle repose sur la marge"]},
    ], 8)

content_slide(prs,
    "Transactions et CA évoluent ensemble — augmenter le CA passe par plus d'achats",
    "04_transactions_par_mois.png",
    [
        {"kpi": "679 k", "kpi_color": "green", "sub": "transactions sur 2 ans",
         "points": ["Moyenne : 28 296 / mois"]},
        {"points": ["Volume et CA parfaitement synchrones",
                    "Aucun effet prix détectable",
                    "→ Pour croître : augmenter la fréquence d'achat"]},
    ], 9)

content_slide(prs,
    "Cat. 0 domine en volume — Cat. 2 génère autant de CA avec bien moins d'unités vendues",
    "05c_qte_vendus_categ.png",
    [
        {"points": ["Cat. 0 = volume fort, prix bas",
                    "Cat. 2 = volume faible, prix élevé"]},
        {"points": ["→ Le CA de Cat. 2 repose sur la marge, pas le volume",
                    "→ Une stratégie de montée en gamme (Cat. 2) est complémentaire à Cat. 0"]},
    ], 10)

content_slide(prs,
    "Top 10 = CA et volume alignés — les flops vendent peu et génèrent un CA négligeable",
    "06b_top_flop_ca_qte.png",
    [
        {"points": ["Top produit : ~95 k€ de CA", "Le top 10 est dominé par Cat. 2"]},
        {"kpi": "< 3 €", "kpi_color": "red", "sub": "CA généré par les flops",
         "points": ["→ Déréférencer les flops : simplifier l'offre sans perte de CA"]},
    ], 11)


# ── Slides 12-14 — Clients ───────────────────────────────────────────────────
content_slide(prs,
    "4 clients BtoB = 7,4% du CA total — perdre un seul impacte directement les résultats",
    "08_ca_btob.png",
    [
        {"kpi": "4", "kpi_color": "red", "sub": "clients BtoB identifiés",
         "points": ["881 k€ générés — 7,4 % du CA total",
                    "→ Concentration critique sur très peu de comptes"]},
        {"points": ["→ Priorité absolue : suivi dédié + contrats cadres",
                    "→ Sécuriser ces comptes avant toute autre action commerciale"]},
    ], 12)

content_slide(prs,
    "Gini 0,45 — concentration modérée, la base BtoC est saine hors les 4 comptes BtoB",
    "09_lorenz.png",
    [
        {"kpi": "48 %", "kpi_color": "neutral", "sub": "du CA sur les top 20 % clients",
         "points": ["Indice de Gini : 0,446 — concentration modérée"]},
        {"points": ["La base BtoC est relativement équilibrée",
                    "Le risque de concentration vient des 4 BtoB, pas du BtoC",
                    "→ Sécuriser le BtoB sans négliger la valeur diffuse du BtoC"]},
    ], 13)

content_slide(prs,
    "Acquisition régulière — mais comprendre quand et pourquoi les clients reviennent est clé",
    "nouveaux_clients.png",
    [
        {"kpi": "8 596", "kpi_color": "green", "sub": "clients uniques sur 2 ans",
         "points": ["Acquisition stable et continue sur la période",
                    "Pas de pic d'acquisition isolé — régularité du recrutement"]},
        {"points": ["→ La base clients est saine et se renouvelle naturellement",
                    "→ L'enjeu : transformer chaque nouveau client en client récurrent"]},
    ], 14)


# ── Slides 15-19 — Corrélations ──────────────────────────────────────────────
content_slide(prs,
    "Genre et catégorie : lien trop faible pour segmenter — l'effort ne sera pas rentabilisé",
    "10_genre_categorie.png",
    [
        {"points": ["Lien statistiquement réel mais de très petite taille",
                    "→ Le genre influence marginalement les préférences d'achat"]},
        {"points": ["→ Segmenter par genre seul : trop peu d'impact sur le CA",
                    "→ À croiser avec l'âge pour affiner le ciblage éditorial"]},
    ], 15)

content_slide(prs,
    "Tranche d'âge et catégorie : lien fort — c'est le seul levier démographique actionnable",
    "14_age_categorie_kruskal.png",
    [
        {"points": ["Lien statistiquement très fort (Kruskal-Wallis, p < 0,001)",
                    "→ L'âge prédit bien la catégorie achetée"]},
        {"points": ["→ Adapter les mises en avant selon la tranche d'âge du client",
                    "→ Résultat le plus actionnable des corrélations analysées",
                    "→ Ciblage éditorial par âge : pertinent et immédiatement opérationnel"]},
    ], 16)

content_slide(prs,
    "Âge et dépenses totales : légère corrélation négative — l'âge seul n'est pas prédicteur de valeur",
    "11_age_ca_total.png",
    [
        {"kpi": "−11 €", "kpi_color": "red", "sub": "de CA par année d'âge supplémentaire",
         "points": ["Corrélation faible — lien réel mais de faible amplitude"]},
        {"points": ["→ L'âge ne suffit pas à prédire la valeur d'un client",
                    "→ Ne pas cibler par âge pour optimiser la valeur client"]},
    ], 17)

content_slide(prs,
    "Âge et fréquence : aucun lien — tous les âges achètent aussi souvent",
    "12_age_frequence.png",
    [
        {"points": ["Aucune relation entre l'âge et la fréquence d'achat",
                    "Un client de 25 ans n'achète pas plus souvent qu'un client de 60 ans"]},
        {"points": ["→ La fréquence ne dépend pas de l'âge",
                    "→ Stratégie : fidéliser tôt, entretenir la relation dans la durée"]},
    ], 18)

content_slide(prs,
    "Âge et panier moyen : aucun lien — tous les âges dépensent autant par achat",
    "13_age_panier_moyen.png",
    [
        {"points": ["Aucune relation entre l'âge et le montant par transaction",
                    "Résultat cohérent avec la fréquence"]},
        {"points": ["Conclusion sur l'âge :",
                    "→ Utile pour le ciblage éditorial (quelle catégorie proposer)",
                    "→ Pas pertinent pour prédire la valeur ou la fréquence d'achat"]},
    ], 19)


# ── Slide 20 — Synthèse ──────────────────────────────────────────────────────
s_syn = prs.slides.add_slide(BLANK)
cy    = header(s_syn, "Synthèse — Ce que les données nous disent", 20)

blocs = [
    ("Activité", GREEN, [
        "CA +11,85 M€ stable sur 2 ans",
        "Cat. 2 : premium sous-exploité",
        "Volume = moteur de croissance",
        "Flops catalogue : CA négligeable",
    ]),
    ("Clients", RGBColor(0x1A, 0x53, 0x76), [
        "4 clients BtoB = 7,4 % du CA",
        "Concentration modérée (Gini 0,45)",
        "Acquisition régulière et stable",
        "Base BtoC saine et équilibrée",
    ]),
    ("Corrélations", RGBColor(0x6C, 0x3A, 0x83), [
        "Âge → fort prédicteur de catégorie",
        "Genre → effet trop faible pour segmenter",
        "Âge → effet faible sur le CA total",
        "Âge → aucun lien avec la fréquence",
    ]),
    ("Signal fort", RGBColor(0x1B, 0x6B, 0x43), [
        "Cibler par âge pour l'offre éditoriale",
        "Sécuriser les 4 comptes BtoB en urgence",
        "Ne pas segmenter par genre",
        "Développer Cat. 2 sur les < 25 ans",
    ]),
]

n_blocs  = len(blocs)
cards_w  = SLIDE_W - 2 * CARD_MARGIN
bloc_w   = (cards_w - Cm(0.3) * (n_blocs - 1)) / n_blocs
bloc_h   = SLIDE_H - cy - Cm(2.0)
HDR_H    = Cm(2.0)

for i, (titre, color, items) in enumerate(blocs):
    x = CARD_MARGIN + i * (bloc_w + Cm(0.3))

    hdr = s_syn.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, cy, bloc_w, HDR_H
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()

    tb_h = s_syn.shapes.add_textbox(x + Cm(0.3), cy + Cm(0.5), bloc_w - Cm(0.6), HDR_H - Cm(0.5))
    tf_h = tb_h.text_frame
    tf_h.word_wrap = True
    r_h  = tf_h.paragraphs[0].add_run()
    r_h.text = titre
    ft(r_h, CARD_FONT, bold=True, color=WHITE)

    box = s_syn.shapes.add_textbox(x, cy + HDR_H, bloc_w, bloc_h - HDR_H)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = BORDER
    box.line.width = Pt(0.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Cm(0.45)

    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(12)
        r = p.add_run()
        r.text = "→  " + item
        ft(r, CARD_FONT)


# ── Slide 21 — Recommandations ───────────────────────────────────────────────
s_reco = prs.slides.add_slide(BLANK)
cy     = header(s_reco, "4 actions prioritaires pour les 6 prochains mois", 21)

recs = [
    (KPI_RED,               "1 — Sécuriser le BtoB",
     "7,4 % du CA sur 4 comptes\n→ Suivi dédié + contrats cadres"),
    (KPI_RED,               "2 — Rationaliser le catalogue",
     "Focus 20 % des références motrices\n→ Déréférencer les flops sans perte de CA"),
    (RGBColor(0xE6,0x7E,0x22), "3 — Développer Cat. 2 sur les < 25 ans",
     "Prix premium + tranche d'âge validée\n→ Seul segment où Cat. 2 est surreprésentée"),
    (RGBColor(0xE6,0x7E,0x22), "4 — Cibler par tranche d'âge",
     "Adapter les mises en avant éditoriales\n→ Seul levier démographique fort (Kruskal p < 0,001)"),
]

n_recs   = len(recs)
cards_w  = SLIDE_W - 2 * CARD_MARGIN
rec_w    = (cards_w - Cm(0.3) * (n_recs - 1)) / n_recs
rec_h    = SLIDE_H - cy - Cm(2.0)
BAND_H   = Cm(2.0)

for i, (color, titre, desc) in enumerate(recs):
    x = CARD_MARGIN + i * (rec_w + Cm(0.3))

    band = s_reco.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, cy, rec_w, BAND_H
    )
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()

    tb_band = s_reco.shapes.add_textbox(x + Cm(0.3), cy + Cm(0.5), rec_w - Cm(0.6), BAND_H - Cm(0.5))
    tf_band = tb_band.text_frame
    tf_band.word_wrap = True
    r_band = tf_band.paragraphs[0].add_run()
    r_band.text = titre
    ft(r_band, CARD_FONT, bold=True, color=WHITE)

    box = s_reco.shapes.add_textbox(x, cy + BAND_H, rec_w, rec_h - BAND_H)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = BORDER
    box.line.width = Pt(0.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Cm(0.45)

    p2 = tf.paragraphs[0]
    r2 = p2.add_run()
    r2.text = desc
    ft(r2, CARD_FONT, color=MID)


# ── Save ─────────────────────────────────────────────────────────────────────
output = "presentation_lapage_v4.pptx"
prs.save(output)
print(f"Sauvegardé : {output} — {len(prs.slides)} slides")
