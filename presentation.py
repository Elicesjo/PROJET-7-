from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

FIG_DIR = Path("figures")
OUT = Path("Lapage_analyse_ecommerce.pptx")

# ── Palette ───────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)   # titres principaux
SLATE      = RGBColor(0x47, 0x55, 0x69)    # texte secondaire
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)   # fond slides
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLUE       = RGBColor(0x25, 0x63, 0xEB)    # accent principal
TEAL       = RGBColor(0x0D, 0x94, 0x88)    # accent secondaire
STEEL      = RGBColor(0x64, 0x74, 0x8B)    # texte corps
DIVIDER    = RGBColor(0xE2, 0xE8, 0xF0)    # lignes de séparation
CAPTION    = RGBColor(0x94, 0xA3, 0xB8)    # légendes
GREEN      = RGBColor(0x05, 0x96, 0x69)
ORANGE     = RGBColor(0xD9, 0x77, 0x06)
RED        = RGBColor(0xDC, 0x26, 0x26)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide_bg(slide, color=LIGHT_GRAY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txt(slide, text, left, top, width, height,
        size=13, color=STEEL, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text      = text
    p.alignment = align
    p.font.size   = Pt(size)
    p.font.color.rgb = color
    p.font.bold   = bold
    p.font.italic = italic


def hline(slide, top, left=Inches(0.55), width=Inches(12.23), color=DIVIDER, thickness=Pt(1)):
    line = slide.shapes.add_shape(1, left, top, width, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def header(slide, title, subtitle=None):
    """Bandeau gauche vertical coloré + titre."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    txt(slide, title, Inches(0.4), Inches(0.22), Inches(12.5), Inches(0.55),
        size=22, color=NAVY, bold=True)
    if subtitle:
        txt(slide, subtitle, Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.35),
            size=12, color=SLATE, italic=True)
    hline(slide, Inches(1.1), left=Inches(0.4), width=Inches(12.5))


def kpi(slide, value, label, note, left, top, width=Inches(2.35), accent=BLUE):
    """Bloc KPI : valeur grande + label + note descriptive."""
    box = slide.shapes.add_shape(1, left, top, width, Inches(1.55))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = DIVIDER
    box.line.width = Pt(1)

    accent_bar = slide.shapes.add_shape(1, left, top, width, Pt(4))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    txt(slide, value, left + Inches(0.15), top + Inches(0.12),
        width - Inches(0.3), Inches(0.55),
        size=22, color=NAVY, bold=True, align=PP_ALIGN.LEFT)
    txt(slide, label, left + Inches(0.15), top + Inches(0.65),
        width - Inches(0.3), Inches(0.3),
        size=11, color=SLATE, bold=True)
    txt(slide, note, left + Inches(0.15), top + Inches(0.97),
        width - Inches(0.3), Inches(0.45),
        size=9, color=CAPTION)


def bullet_box(slide, title, bullets, left, top, width, height, accent=BLUE):
    txt(slide, title, left, top, width, Inches(0.3),
        size=12, color=accent, bold=True)
    y = top + Inches(0.35)
    for b in bullets:
        txt(slide, f"→  {b}", left, y, width, Inches(0.35), size=11, color=STEEL)
        y += Inches(0.38)


def fig(slide, path, left, top, width, height):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), left, top, width, height)


# ══════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, NAVY)

bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

txt(s, "LAPAGE", Inches(0.9), Inches(1.8), Inches(11), Inches(0.6),
    size=13, color=CAPTION, bold=True, align=PP_ALIGN.LEFT)

txt(s, "Analyse de l'activité\ne-commerce", Inches(0.9), Inches(2.35), Inches(10), Inches(1.6),
    size=36, color=WHITE, bold=True)

line = s.shapes.add_shape(1, Inches(0.9), Inches(4.05), Inches(2.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = BLUE
line.line.fill.background()

txt(s, "Bilan 2 ans d'e-commerce  ·  Mars 2021 – Février 2023",
    Inches(0.9), Inches(4.25), Inches(11), Inches(0.4),
    size=13, color=CAPTION)

txt(s, "Présentation CODIR  ·  Mars 2023  ·  Direction commerciale",
    Inches(0.9), Inches(6.8), Inches(11), Inches(0.35),
    size=10, color=CAPTION)

# ══════════════════════════════════════════════════════════
# SLIDE 2 — Chiffres clés
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "Chiffres clés", "Vue d'ensemble · Mars 2021 – Février 2023")

kpi(s, "11,8 M€",   "Chiffre d'affaires",      "Sur 24 mois",                  Inches(0.4),  Inches(1.35), accent=BLUE)
kpi(s, "678 284",   "Transactions",             "Lignes de vente nettes",       Inches(2.86), Inches(1.35), accent=TEAL)
kpi(s, "8 600",     "Clients actifs",           "Ayant passé ≥ 1 commande",     Inches(5.32), Inches(1.35), accent=GREEN)
kpi(s, "493 k€",    "CA moyen / mois",          "Sur la période",               Inches(7.78), Inches(1.35), accent=ORANGE)
kpi(s, "34,58 €",   "Valeur moy. commande",     "Ce qu'un client dépense / visite", Inches(10.24), Inches(1.35), accent=STEEL)

txt(s, "17,45 €  ·  Prix moyen par article",
    Inches(10.24), Inches(3.0), Inches(2.9), Inches(0.3),
    size=10, color=CAPTION)

hline(s, Inches(3.2), left=Inches(0.4), width=Inches(12.5))

txt(s,
    "La plateforme e-commerce s'est imposée comme un canal majeur dès sa première année. "
    "Avec 8 600 clients actifs et près de 680 000 transactions, l'activité est soutenue et régulière. "
    "La valeur d'une commande (34,58 €) correspond à environ 2 articles au prix moyen (17,45 €).",
    Inches(0.4), Inches(3.35), Inches(12.5), Inches(0.9), size=12, color=STEEL)

txt(s, "Source : base transactions e-commerce Lapage  ·  Données nettoyées (lignes test et doublons exclus)",
    Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3), size=9, color=CAPTION)

# ══════════════════════════════════════════════════════════
# SLIDE 3 — CA mensuel
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "Évolution du chiffre d'affaires", "CA mensuel avec moyenne mobile 3 mois")

fig(s, FIG_DIR / "01_ca_mensuel.png", Inches(0.4), Inches(1.25), Inches(8.8), Inches(4.2))

hline(s, Inches(1.25), left=Inches(9.4), width=Inches(3.5), color=DIVIDER)
bullet_box(s, "Points clés", [
    "Mois record : février 2022",
    "Mois creux : octobre 2021",
    "Tendance stable sur 24 mois",
    "Pas de saisonnalité marquée",
], Inches(9.4), Inches(1.35), Inches(3.5), Inches(3.0))

hline(s, Inches(5.6), left=Inches(0.4))
txt(s,
    "La moyenne mobile (ligne rouge) confirme une activité sans déclin structurel. "
    "L'absence de saisonnalité forte suggère une demande continue, indépendante des périodes clés.",
    Inches(0.4), Inches(5.72), Inches(12.5), Inches(0.7), size=11, color=STEEL)

# ══════════════════════════════════════════════════════════
# SLIDE 4 — CA par catégorie
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "Répartition du CA par catégorie", "Structure des ventes et évolution mensuelle")

fig(s, FIG_DIR / "07_repartition_categorie.png", Inches(0.4), Inches(1.25), Inches(7.5), Inches(3.6))
fig(s, FIG_DIR / "02_ca_par_categorie.png",      Inches(0.4), Inches(4.95), Inches(7.5), Inches(2.3))

hline(s, Inches(1.25), left=Inches(8.1), width=Inches(4.8), color=DIVIDER)
bullet_box(s, "Lecture", [
    "Catégorie 1 : 39% du CA (leader)",
    "Catégorie 0 : 37% du CA",
    "Catégorie 2 : 23% du CA",
    "Équilibre stable sur 24 mois",
    "Cat. 2 : prix unitaires élevés (jusqu'à 300 €)",
], Inches(8.1), Inches(1.35), Inches(4.8), Inches(3.5))

hline(s, Inches(5.0), left=Inches(8.1), width=Inches(4.8), color=DIVIDER)
txt(s,
    "Bien que minoritaire en volume, la catégorie 2 "
    "représente 23% du CA grâce à ses prix premium.",
    Inches(8.1), Inches(5.15), Inches(4.8), Inches(0.8), size=11, color=STEEL)

# ══════════════════════════════════════════════════════════
# SLIDE 5 — Activité clients & transactions
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "Activité clients & transactions", "Évolution mensuelle des volumes")

fig(s, FIG_DIR / "03_clients_par_mois.png",    Inches(0.4),  Inches(1.25), Inches(6.1), Inches(2.6))
fig(s, FIG_DIR / "04_transactions_par_mois.png", Inches(6.7), Inches(1.25), Inches(6.2), Inches(2.6))
fig(s, FIG_DIR / "05_references_par_mois.png", Inches(0.4),  Inches(4.05), Inches(6.1), Inches(2.6))

hline(s, Inches(4.05), left=Inches(6.7), width=Inches(6.2), color=DIVIDER)
bullet_box(s, "Synthèse mensuelle", [
    "~5 700 clients uniques / mois",
    "~28 000 transactions / mois",
    "~2 450 références distinctes vendues",
    "Base clients stable, pas d'érosion",
], Inches(6.7), Inches(4.2), Inches(6.2), Inches(2.5))

# ══════════════════════════════════════════════════════════
# SLIDE 6 — Top / Flop
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "Zoom références — Top & Flop produits", "Classement par chiffre d'affaires généré")

fig(s, FIG_DIR / "06_top_flop_produits.png", Inches(0.4), Inches(1.25), Inches(8.8), Inches(4.5))

hline(s, Inches(1.25), left=Inches(9.4), width=Inches(3.5), color=DIVIDER)
bullet_box(s, "Enseignements", [
    "Top produit (2_159) : 95 k€ de CA",
    "Le top 10 est dominé par la cat. 2",
    "Flops : < 3 € de CA, 1–2 ventes",
], Inches(9.4), Inches(1.35), Inches(3.5), Inches(2.2))

hline(s, Inches(3.7), left=Inches(9.4), width=Inches(3.5), color=DIVIDER)
txt(s, "Action recommandée",
    Inches(9.4), Inches(3.85), Inches(3.5), Inches(0.3),
    size=11, color=BLUE, bold=True)
txt(s,
    "Retirer les références flop du catalogue "
    "pour simplifier l'offre et réduire les coûts de gestion.",
    Inches(9.4), Inches(4.2), Inches(3.5), Inches(0.9), size=11, color=STEEL)

hline(s, Inches(5.8))
txt(s, "Source : calcul sur 678 284 transactions nettes",
    Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.3), size=9, color=CAPTION)

# ══════════════════════════════════════════════════════════
# SLIDE 7 — Lorenz + BtoB
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "Concentration du CA", "Courbe de Lorenz & segmentation clients à fort volume")

fig(s, FIG_DIR / "09_lorenz.png", Inches(0.4),  Inches(1.25), Inches(5.5), Inches(5.5))
fig(s, FIG_DIR / "08_ca_btob.png", Inches(6.1), Inches(1.25), Inches(4.2), Inches(5.5))

hline(s, Inches(1.25), left=Inches(10.5), width=Inches(2.5), color=DIVIDER)
txt(s, "Lecture", Inches(10.5), Inches(1.35), Inches(2.5), Inches(0.3),
    size=11, color=BLUE, bold=True)
txt(s,
    "Gini = 0,446\n\n"
    "Les 20% des meilleurs clients génèrent 48% du CA.\n\n"
    "4 clients à très fort volume (proxy BtoB) concentrent 7,4% du CA.",
    Inches(10.5), Inches(1.75), Inches(2.5), Inches(3.0), size=11, color=STEEL)

hline(s, Inches(5.8))
txt(s,
    "Proxy BtoB : clients dont le CA dépasse moyenne + 2 écarts-types (seuil : 11 719 €)  "
    "·  Concentration modérée : Gini < 0,5",
    Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.3), size=9, color=CAPTION)

# ══════════════════════════════════════════════════════════
# SLIDE 8 — Genre & Catégorie
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "Comportement clients — Genre & Catégorie", "Test d'indépendance Chi²")

fig(s, FIG_DIR / "10_genre_categorie.png", Inches(0.4), Inches(1.25), Inches(8.2), Inches(4.5))

hline(s, Inches(1.25), left=Inches(8.8), width=Inches(4.1), color=DIVIDER)
txt(s, "Résultat statistique", Inches(8.8), Inches(1.35), Inches(4.1), Inches(0.3),
    size=11, color=BLUE, bold=True)
txt(s,
    "Test Chi²  ·  p < 0,05\n→ Lien statistiquement significatif",
    Inches(8.8), Inches(1.75), Inches(4.1), Inches(0.7), size=11, color=STEEL)

hline(s, Inches(2.6), left=Inches(8.8), width=Inches(4.1), color=DIVIDER)
txt(s, "En pratique", Inches(8.8), Inches(2.7), Inches(4.1), Inches(0.3),
    size=11, color=TEAL, bold=True)
txt(s,
    "Femmes : 60,9% cat. 0 — 34,0% cat. 1\n"
    "Hommes : 61,4% cat. 0 — 32,9% cat. 1\n\n"
    "Les différences sont minimes. "
    "Hommes et femmes achètent les mêmes catégories.",
    Inches(8.8), Inches(3.1), Inches(4.1), Inches(1.5), size=11, color=STEEL)

hline(s, Inches(4.75), left=Inches(8.8), width=Inches(4.1), color=DIVIDER)
txt(s, "Recommandation",
    Inches(8.8), Inches(4.85), Inches(4.1), Inches(0.3),
    size=11, color=ORANGE, bold=True)
txt(s,
    "Pas de segmentation marketing par genre sur les catégories.",
    Inches(8.8), Inches(5.22), Inches(4.1), Inches(0.5), size=11, color=STEEL)

# ══════════════════════════════════════════════════════════
# SLIDE 9 — Âge & Achats
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "Comportement clients — Âge & Achats", "Corrélations de Pearson par client")

fig(s, FIG_DIR / "11_age_ca_total.png",    Inches(0.4), Inches(1.25), Inches(6.1), Inches(2.8))
fig(s, FIG_DIR / "13_age_panier_moyen.png", Inches(6.7), Inches(1.25), Inches(6.2), Inches(2.8))
fig(s, FIG_DIR / "12_age_frequence.png",   Inches(0.4), Inches(4.2),  Inches(6.1), Inches(2.8))

hline(s, Inches(4.2), left=Inches(6.7), width=Inches(6.2), color=DIVIDER)
bullet_box(s, "Synthèse", [
    "Panier moyen : corrélation forte r = −0,51",
    "  → Les jeunes dépensent plus par achat",
    "CA total : corrélation très faible r = −0,04",
    "Fréquence : aucun lien (p = 0,53)",
    "  → Tous les âges achètent autant",
], Inches(6.7), Inches(4.35), Inches(6.2), Inches(2.8))

# ══════════════════════════════════════════════════════════
# SLIDE 10 — Âge & Catégorie
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header(s, "Comportement clients — Âge & Catégorie", "Test Chi²  ·  Résultat très significatif (p < 0,001)")

fig(s, FIG_DIR / "14_age_categorie.png", Inches(0.4), Inches(1.25), Inches(8.5), Inches(4.8))

hline(s, Inches(1.25), left=Inches(9.1), width=Inches(3.8), color=DIVIDER)
txt(s, "Par tranche d'âge", Inches(9.1), Inches(1.35), Inches(3.8), Inches(0.3),
    size=11, color=BLUE, bold=True)

tranches = [
    ("< 25 ans",  "42% cat. 2 (livres premium)",      RED),
    ("25–34 ans", "54% cat. 0",                        STEEL),
    ("35–49 ans", "76% cat. 0 (très concentré)",       STEEL),
    ("50–64 ans", "48% cat. 1",                        STEEL),
    ("65+ ans",   "56% cat. 1",                        TEAL),
]
y = Inches(1.75)
for tranche, detail, color in tranches:
    txt(s, tranche, Inches(9.1), y, Inches(1.4), Inches(0.32),
        size=10, color=color, bold=True)
    txt(s, detail,  Inches(10.6), y, Inches(2.3), Inches(0.32),
        size=10, color=STEEL)
    y += Inches(0.4)

hline(s, Inches(3.95), left=Inches(9.1), width=Inches(3.8), color=DIVIDER)
txt(s, "Recommandation",
    Inches(9.1), Inches(4.1), Inches(3.8), Inches(0.3),
    size=11, color=ORANGE, bold=True)
txt(s,
    "Cibler les < 25 ans sur la catégorie 2 "
    "(coffrets, livres premium). Ce segment "
    "a le panier moyen le plus élevé.",
    Inches(9.1), Inches(4.5), Inches(3.8), Inches(1.0), size=11, color=STEEL)

# ══════════════════════════════════════════════════════════
# SLIDE 11 — Recommandations
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, WHITE)

bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

txt(s, "Recommandations stratégiques",
    Inches(0.4), Inches(0.22), Inches(12.5), Inches(0.55),
    size=22, color=NAVY, bold=True)
txt(s, "Synthèse des actions prioritaires",
    Inches(0.4), Inches(0.77), Inches(12.5), Inches(0.3),
    size=12, color=SLATE, italic=True)
hline(s, Inches(1.1), left=Inches(0.4), width=Inches(12.5))

recos = [
    (BLUE,   "01",  "Fidéliser les clients à fort volume (BtoB)",
     "4 clients représentent 7,4% du CA. Un programme dédié (compte pro, remises volume) réduit le risque de perte."),
    (TEAL,   "02",  "Cibler les < 25 ans sur la catégorie 2",
     "Ce segment achète 42% de livres premium avec le panier moyen le plus élevé. Des offres cadeaux/coffrets sont pertinentes."),
    (ORANGE, "03",  "Rationaliser le catalogue",
     "Les références flop (< 3 € de CA, 1–2 ventes) n'apportent pas de valeur. Leur retrait simplifie l'expérience client."),
    (STEEL,  "04",  "Pas de segmentation par genre",
     "Les comportements d'achat sont quasi-identiques entre hommes et femmes. Réorienter ces ressources marketing."),
]

for i, (color, num, titre, detail) in enumerate(recos):
    top = Inches(1.3 + i * 1.4)

    badge = s.shapes.add_shape(1, Inches(0.4), top + Inches(0.1), Inches(0.45), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    txt(s, num, Inches(0.4), top + Inches(0.1), Inches(0.45), Inches(0.45),
        size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    txt(s, titre,  Inches(1.05), top + Inches(0.08), Inches(11.5), Inches(0.35),
        size=13, color=NAVY, bold=True)
    txt(s, detail, Inches(1.05), top + Inches(0.45), Inches(11.5), Inches(0.6),
        size=11, color=STEEL)

    if i < 3:
        hline(s, top + Inches(1.25), left=Inches(0.4), color=DIVIDER)

txt(s, "Lapage  ·  Analyse e-commerce  ·  Mars 2023",
    Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3),
    size=9, color=CAPTION, align=PP_ALIGN.RIGHT)

# ── Sauvegarde ────────────────────────────────────────────
prs.save(OUT)
print(f"Présentation sauvegardée → {OUT}  ({len(prs.slides)} slides)")
